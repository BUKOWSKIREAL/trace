"""
PptxHandler — .pptx 幻灯片的"幻灯片级" diff / 摘要
=====================================================
通过 python-pptx 把每张幻灯片上的文本框与表格单元格按出现顺序拼成文本流，
为每张幻灯片加一个 "=== Slide N ===" 头部行，便于按幻灯片维度阅读差异。

# AI生成 (Claude 4.7 协助)
# 人工修正：
#   - 头部行 "=== Slide N ===" 是语义锚点，让多张幻灯片的 diff 不会糊在一起；
#   - shape 既可能 has_text_frame 也可能 has_table，要分别处理；
#     图片 / 组合形状没有文本，跳过即可——不要让它进入文本流；
#   - python-pptx 解析坏包会抛 zipfile.BadZipFile / KeyError / PackageNotFoundError
#     等多种异常，一律 catch 成 meta 行，不让异常冒到 UI 层。
#   - 追踪盲区修复：补读演讲者备注（notes_slide）+ 递归组合形状（group）里的
#     文本框——原版改备注/组合文字会被误报"幻灯片级一致"。
"""
from __future__ import annotations

import difflib
import logging
from io import BytesIO

from core.handlers.base import DiffLine, FileHandler, HandlerRegistry

logger = logging.getLogger("trace.handlers.pptx")


class PptxHandler(FileHandler):
    """
    .pptx 文件 handler，按幻灯片为语义单元做行级 diff。
    """

    extensions = [".pptx"]

    def extract_text(self, blob: bytes) -> str | None:
        """
        把每张幻灯片的文本框 / 表格 / 备注按顺序拼接，每张幻灯片前加 "=== Slide N ===" 头。
        组合形状（group）里的文字会被递归取出。
        损坏的 .pptx 返回 None——调用方据此把它当成二进制处理。
        """
        try:
            from pptx import Presentation  # 延迟导入，避免无 pptx 时整模块崩
        except ImportError:
            logger.warning("python-pptx 未安装，PptxHandler 退化为不可解析")
            return None

        try:
            prs = Presentation(BytesIO(blob))
        except Exception as e:
            logger.warning("PptxHandler 解析失败：%s", e)
            return None

        lines: list[str] = []
        for i, slide in enumerate(prs.slides):
            lines.append(f"=== Slide {i + 1} ===")
            for shape in slide.shapes:
                # 递归收集形状文字——含组合形状（group）里嵌套的文本框 / 表格
                self._collect_shape_text(shape, lines)
            # 备注页（演讲者备注）——用户实际会编辑，原版完全漏掉
            if slide.has_notes_slide:
                try:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes and notes.strip():
                        for para in notes.splitlines():
                            lines.append(f"[备注] {para}")
                except Exception:
                    # 个别 pptx 的 notes_slide 结构异常，跳过即可
                    pass
        return "\n".join(lines)

    def _collect_shape_text(self, shape, lines: list[str]) -> None:
        """
        递归收集单个 shape 的文本，把结果 append 到 lines。

        # 人工修正（追踪盲区修复）：原版只处理顶层 shape 的 text_frame/table，
        # 组合形状（GROUP）里嵌套的文本框文字被整体跳过——改了组合里的文字
        # 会显示"幻灯片级一致"。这里对 group 递归下钻。
        """
        # 组合形状：递归进子形状。能走到这里说明 pptx 已装好，import 命中缓存。
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            is_group = getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_group = False
        if is_group:
            try:
                for sub in shape.shapes:
                    self._collect_shape_text(sub, lines)
            except Exception:
                pass
            return
        # 文本框（含 title、占位符、普通 textbox）
        if shape.has_text_frame:
            # text_frame.text 已含段落间的 '\n'；按段落展开保证细粒度 diff
            for para in shape.text_frame.paragraphs:
                lines.append(para.text)
        # 表格——按行 / 单元格展开，让单元格变更可见
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    lines.append(cell.text)
        # 其他类型（图片 / 图表）没有可比对文本，跳过

    def describe_change(self, old_blob: bytes, new_blob: bytes) -> str:
        """
        返回 "增加 X 张幻灯片 / 删除 Y 张 / 修改 Z 张"。
        # 人工注释：以 "=== Slide N ===" 头切分文本，按幻灯片整体内容比对。
        # 用 SequenceMatcher 区分 insert / delete / replace 三类操作，
        # 而不是简单数行数——这样能反映"换了一张幻灯片"这种语义。
        """
        old_slides = self._split_slides(old_blob)
        new_slides = self._split_slides(new_blob)
        if old_slides is None or new_slides is None:
            return "[文件解析失败：无法读取 .pptx 幻灯片]"

        sm = difflib.SequenceMatcher(a=old_slides, b=new_slides, autojunk=False)
        added = removed = modified = 0
        for op, i1, i2, j1, j2 in sm.get_opcodes():
            if op == "insert":
                added += j2 - j1
            elif op == "delete":
                removed += i2 - i1
            elif op == "replace":
                modified += max(i2 - i1, j2 - j1)

        if added == 0 and removed == 0 and modified == 0:
            return "幻灯片级一致（仅样式或非文本元素差异）"
        return f"增加 {added} 张幻灯片 / 删除 {removed} 张 / 修改 {modified} 张"

    def render_diff(self, old_blob: bytes, new_blob: bytes) -> list[DiffLine]:
        """
        以"段落 / 单元格 / 头部"为最小行做 ndiff——
        头部行 "=== Slide N ===" 会自然出现在 diff 里，UI 据此分段。
        坏 .pptx 返回单条 meta，让 UI 显式告诉用户"解析失败"。
        """
        if not old_blob and not new_blob:
            return []

        if not old_blob:
            new_text = self.extract_text(new_blob)
            if new_text is None:
                return [("meta", "[文件解析失败：无法读取 .pptx 幻灯片]")]
            return self._text_as_diff(new_text, "added")

        if not new_blob:
            old_text = self.extract_text(old_blob)
            if old_text is None:
                return [("meta", "[文件解析失败：无法读取 .pptx 幻灯片]")]
            return self._text_as_diff(old_text, "removed")

        old_text = self.extract_text(old_blob)
        new_text = self.extract_text(new_blob)
        if old_text is None or new_text is None:
            return [("meta", "[文件解析失败：无法读取 .pptx 幻灯片]")]

        old_lines = old_text.splitlines()
        new_lines = new_text.splitlines()
        result: list[DiffLine] = []
        for line in difflib.ndiff(old_lines, new_lines):
            tag2 = line[:2]
            content = line[2:]
            if tag2 == "  ":
                result.append(("normal", content))
            elif tag2 == "+ ":
                result.append(("added", content))
            elif tag2 == "- ":
                result.append(("removed", content))
            # '? ' 是 difflib 的位置提示，跳过

        return result

    # ---------- 内部辅助 ----------

    def _split_slides(self, blob: bytes) -> list[str] | None:
        """
        按 "=== Slide N ===" 头切分整篇文本，返回每张幻灯片的内容块列表。
        # 人工注释：把"幻灯片"当成 diff 的最小单元——增删一整张幻灯片
        # 不应被错算成"一堆行的新增/删除"，应该是"+1 张 / -1 张"。
        """
        text = self.extract_text(blob)
        if text is None:
            return None
        slides: list[str] = []
        buf: list[str] = []
        for line in text.splitlines():
            if line.startswith("=== Slide ") and line.endswith(" ==="):
                if buf:
                    slides.append("\n".join(buf))
                    buf = []
            buf.append(line)
        if buf:
            slides.append("\n".join(buf))
        return slides


# 模块导入时自动注册
HandlerRegistry.register(PptxHandler())
