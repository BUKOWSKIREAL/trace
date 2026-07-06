"""
Binary-ish handler lifecycle tests.

These handlers must render useful output for file creation and deletion.
The repository represents "no previous version" and "deleted version" as an
empty blob, so handlers must not parse b"" as a corrupt Office/PDF/image file.
"""
import sys
import unittest
from io import BytesIO
from pathlib import Path

ROOT = Path(__file__).parent.parent / "code"
sys.path.insert(0, str(ROOT))

import fitz  # noqa: E402
from docx import Document  # noqa: E402
from openpyxl import Workbook  # noqa: E402
from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402

from core.handlers.binary_handler import BinaryHandler  # noqa: E402
from core.handlers.docx_handler import DocxHandler  # noqa: E402
from core.handlers.image_handler import ImageHandler  # noqa: E402
from core.handlers.pdf_handler import PdfHandler  # noqa: E402
from core.handlers.pptx_handler import PptxHandler  # noqa: E402
from core.handlers.xlsx_handler import XlsxHandler  # noqa: E402


def make_docx(text: str) -> bytes:
    doc = Document()
    doc.add_paragraph(text)
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_pptx(text: str) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if slide.shapes.title is not None:
        slide.shapes.title.text = text
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def make_xlsx(text: str) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws["A1"] = text
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def make_pdf(text: str) -> bytes:
    doc = fitz.open()
    try:
        page = doc.new_page()
        page.insert_text((72, 72), text)
        return doc.tobytes()
    finally:
        doc.close()


def make_png() -> bytes:
    img = Image.new("RGB", (4, 4), color="red")
    buf = BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


class TestBinaryFileLifecycleHandlers(unittest.TestCase):
    def test_new_supported_binary_files_do_not_report_parse_failure(self):
        cases = [
            (DocxHandler(), make_docx("new docx text"), "new docx text"),
            (PptxHandler(), make_pptx("new pptx text"), "new pptx text"),
            (XlsxHandler(), make_xlsx("new xlsx text"), "new xlsx text"),
            (PdfHandler(), make_pdf("new pdf text"), "new pdf text"),
            (ImageHandler(), make_png(), "SHA-256"),
            (BinaryHandler(), b"raw bytes", "二进制文件"),
        ]
        for handler, blob, expected in cases:
            with self.subTest(handler=type(handler).__name__):
                out = handler.render_diff(b"", blob)
                joined = "\n".join(text for _, text in out)
                self.assertNotIn("解析失败", joined)
                self.assertIn(expected, joined)
                if type(handler).__name__ in {"DocxHandler", "PptxHandler", "XlsxHandler", "PdfHandler"}:
                    tags = {tag for tag, _ in out}
                    self.assertIn("added", tags)

    def test_deleted_supported_binary_files_do_not_report_parse_failure(self):
        cases = [
            (DocxHandler(), make_docx("old docx text"), "old docx text"),
            (PptxHandler(), make_pptx("old pptx text"), "old pptx text"),
            (XlsxHandler(), make_xlsx("old xlsx text"), "old xlsx text"),
            (PdfHandler(), make_pdf("old pdf text"), "old pdf text"),
            (ImageHandler(), make_png(), "SHA-256"),
            (BinaryHandler(), b"raw bytes", "二进制文件"),
        ]
        for handler, blob, expected in cases:
            with self.subTest(handler=type(handler).__name__):
                out = handler.render_diff(blob, b"")
                joined = "\n".join(text for _, text in out)
                self.assertNotIn("解析失败", joined)
                self.assertIn(expected, joined)
                if type(handler).__name__ in {"DocxHandler", "PptxHandler", "XlsxHandler", "PdfHandler"}:
                    tags = {tag for tag, _ in out}
                    self.assertIn("removed", tags)


if __name__ == "__main__":
    unittest.main(verbosity=2)
