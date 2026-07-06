"""
PptxHandler 单元测试
=======================
覆盖 PptxHandler 的 extract_text / describe_change / render_diff / 鲁棒性 / 注册路由。

# 人工编写
"""
import sys
import unittest
from io import BytesIO
from pathlib import Path

# 让 tests/ 也能跑（对齐 tests/test_handlers.py 的写法）
ROOT = Path(__file__).parent.parent / "code"
sys.path.insert(0, str(ROOT))

from core.handlers import HandlerRegistry  # noqa: E402
from core.handlers.pptx_handler import PptxHandler  # noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402


def _build_pptx(slides_spec: list[dict]) -> bytes:
    """
    用 python-pptx 在内存里构造一个 .pptx blob。
    slides_spec 示例：
        [{"title": "Hello", "boxes": ["line a", "line b"], "table": [["A","B"],["C","D"]]}]
    """
    prs = Presentation()
    blank_layout = prs.slide_layouts[5]  # Title only layout
    for spec in slides_spec:
        slide = prs.slides.add_slide(blank_layout)
        title = spec.get("title")
        if title is not None and slide.shapes.title is not None:
            slide.shapes.title.text = title
        boxes = spec.get("boxes", [])
        for idx, text in enumerate(boxes):
            tb = slide.shapes.add_textbox(
                Inches(1), Inches(1.5 + idx * 0.6), Inches(6), Inches(0.5)
            )
            tb.text_frame.text = text
        table_data = spec.get("table")
        if table_data:
            rows = len(table_data)
            cols = len(table_data[0])
            table_shape = slide.shapes.add_table(
                rows, cols, Inches(1), Inches(4), Inches(6), Inches(1)
            )
            table = table_shape.table
            for r, row in enumerate(table_data):
                for c, val in enumerate(row):
                    table.cell(r, c).text = val
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestExtractText(unittest.TestCase):
    def setUp(self):
        self.h = PptxHandler()

    def test_extract_text_contains_slide_header(self):
        blob = _build_pptx([{"title": "Hello"}])
        text = self.h.extract_text(blob)
        self.assertIsNotNone(text)
        self.assertIn("=== Slide 1 ===", text)
        self.assertIn("Hello", text)

    def test_extract_text_multi_slide_headers(self):
        blob = _build_pptx([{"title": "S1"}, {"title": "S2"}, {"title": "S3"}])
        text = self.h.extract_text(blob)
        self.assertIn("=== Slide 1 ===", text)
        self.assertIn("=== Slide 2 ===", text)
        self.assertIn("=== Slide 3 ===", text)
        self.assertIn("S1", text)
        self.assertIn("S2", text)
        self.assertIn("S3", text)

    def test_extract_text_includes_textbox_and_table_cells(self):
        blob = _build_pptx(
            [
                {
                    "title": "T",
                    "boxes": ["box-a", "box-b"],
                    "table": [["cellX", "cellY"], ["cellZ", "cellW"]],
                }
            ]
        )
        text = self.h.extract_text(blob)
        for token in ("box-a", "box-b", "cellX", "cellY", "cellZ", "cellW"):
            self.assertIn(token, text, f"missing token {token}")


class TestDescribeChange(unittest.TestCase):
    def setUp(self):
        self.h = PptxHandler()

    def test_describe_change_add_slide(self):
        old = _build_pptx([{"title": "A"}])
        new = _build_pptx([{"title": "A"}, {"title": "B"}])
        desc = self.h.describe_change(old, new)
        # 增加一张幻灯片
        self.assertIn("增加", desc)
        self.assertIn("1", desc)

    def test_describe_change_remove_slide(self):
        old = _build_pptx([{"title": "A"}, {"title": "B"}])
        new = _build_pptx([{"title": "A"}])
        desc = self.h.describe_change(old, new)
        self.assertIn("删除", desc)
        self.assertIn("1", desc)

    def test_describe_change_modify_slide(self):
        old = _build_pptx([{"title": "A", "boxes": ["old text"]}])
        new = _build_pptx([{"title": "A", "boxes": ["new text"]}])
        desc = self.h.describe_change(old, new)
        # 同样张数，应当是修改
        self.assertIn("修改", desc)

    def test_describe_change_identical_textual_content(self):
        old = _build_pptx([{"title": "Same"}])
        new = _build_pptx([{"title": "Same"}])
        desc = self.h.describe_change(old, new)
        # 文本完全相同——可能命中"幻灯片级一致"，也可能命中 0/0/0
        self.assertIn("一致", desc)


class TestRenderDiff(unittest.TestCase):
    def setUp(self):
        self.h = PptxHandler()

    def test_render_diff_tags_present(self):
        old = _build_pptx([{"title": "A", "boxes": ["alpha"]}])
        new = _build_pptx([{"title": "A", "boxes": ["beta"]}])
        out = self.h.render_diff(old, new)
        tags = {t for t, _ in out}
        self.assertIn("added", tags)
        self.assertIn("removed", tags)
        # 不应出现 difflib 的 '? ' hint
        self.assertNotIn("?", tags)

    def test_render_diff_no_change_only_normal(self):
        blob = _build_pptx([{"title": "X"}])
        out = self.h.render_diff(blob, blob)
        tags = {t for t, _ in out}
        # 同 blob 进出——应当全是 normal
        self.assertEqual(tags, {"normal"})

    def test_render_diff_contents_visible(self):
        old = _build_pptx([{"title": "A", "boxes": ["old-content"]}])
        new = _build_pptx([{"title": "A", "boxes": ["new-content"]}])
        out = self.h.render_diff(old, new)
        joined = " ".join(c for _, c in out)
        self.assertIn("old-content", joined)
        self.assertIn("new-content", joined)


class TestRegistry(unittest.TestCase):
    def test_pptx_routes_to_pptx_handler(self):
        h = HandlerRegistry.for_path(Path("deck.pptx"))
        self.assertIsInstance(h, PptxHandler)

    def test_pptx_extension_case_insensitive(self):
        h = HandlerRegistry.for_path(Path("DECK.PPTX"))
        self.assertIsInstance(h, PptxHandler)


class TestRobustness(unittest.TestCase):
    def setUp(self):
        self.h = PptxHandler()

    def test_corrupt_blob_extract_returns_none(self):
        # 完全不是 zip 的垃圾字节——不应抛
        self.assertIsNone(self.h.extract_text(b"not a real pptx"))

    def test_corrupt_blob_describe_returns_meta(self):
        desc = self.h.describe_change(b"not a pptx", b"also garbage")
        self.assertIn("解析失败", desc)

    def test_corrupt_blob_render_returns_meta_line(self):
        out = self.h.render_diff(b"trash-1", b"trash-2")
        self.assertEqual(len(out), 1)
        tag, content = out[0]
        self.assertEqual(tag, "meta")
        self.assertIn("解析失败", content)


class TestNotesTracking(unittest.TestCase):
    """追踪盲区修复：演讲者备注必须可见（原版完全漏读 notes_slide）。"""

    def setUp(self):
        self.h = PptxHandler()

    @staticmethod
    def _pptx_with_notes(notes: str) -> bytes:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        if slide.shapes.title is not None:
            slide.shapes.title.text = "标题不变"
        slide.notes_slide.notes_text_frame.text = notes
        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def test_notes_text_extracted(self):
        out = self.h.extract_text(self._pptx_with_notes("演讲者备注内容"))
        self.assertIn("演讲者备注内容", out)
        self.assertIn("备注", out)  # 带标记前缀

    def test_notes_change_is_visible(self):
        # 正文不变、只改备注——原版会误报"幻灯片级一致"
        old = self._pptx_with_notes("备注第一版")
        new = self._pptx_with_notes("备注第二版")
        desc = self.h.describe_change(old, new)
        self.assertNotIn("一致", desc)
        contents = " ".join(c for _, c in self.h.render_diff(old, new))
        self.assertIn("备注第一版", contents)
        self.assertIn("备注第二版", contents)


if __name__ == "__main__":
    unittest.main(verbosity=2)
