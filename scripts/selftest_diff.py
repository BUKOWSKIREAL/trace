#!/usr/bin/env python3
"""
selftest_diff.py — 五种二进制格式 diff 的确定性自测
=======================================================
一条命令、不依赖后台 daemon、不依赖 GUI，让用户亲眼看到
docx / pptx / xlsx / pdf / png 五种二进制文件的"具体文字改动"。

做法（完全复刻 views/diff_view.py 的真实 DiffView 管线，只是把
tk.Text 换成终端 ANSI 着色）：
    1. 在系统临时目录建一个干净 workspace
    2. 生成 5 个文件的 v1（python-docx / python-pptx / openpyxl /
       PyMuPDF(fitz) / Pillow），用 Repository commit 一次 → commit #1
    3. 覆盖同名文件生成 v2（制造能明显看出的文字变化），再 commit
       一次 → commit #2
    4. 对 commit#1 → #2：用 _load_manifest 取两份 manifest，对每个
       hash 不同的文件 storage.get 取字节，HandlerRegistry.for_path
       分发 handler，打印 describe_change 摘要 + render_diff 逐行
    5. 结尾打印明确的通过 / 失败提示；finally 里清理临时目录

# 人工编写（自测脚本，独立于项目其它源码，仅新增此文件）
"""
from __future__ import annotations

import shutil
import sys
import tempfile
from pathlib import Path

# === 把项目的 code/ 加入 sys.path（绝对路径）===
# 因为打包成 .app 后 cwd 任意，必须用脚本自身位置推算 code/ 的绝对路径，
# 不能依赖相对 cwd。脚本在 <repo>/scripts/，code 在 <repo>/code/。
_REPO_ROOT = Path(__file__).resolve().parent.parent
_CODE_DIR = _REPO_ROOT / "code"
for _p in (_CODE_DIR):
    if _p.is_dir() and str(_p) not in sys.path:
        sys.path.insert(0, str(_p))

# 项目内部 import（必须在 sys.path 设置之后）
from core.handlers import HandlerRegistry  # noqa: E402  触发全部 handler 注册
from core.repository import Repository  # noqa: E402
from models.agent import AgentAttribution  # noqa: E402
from models.change import Change  # noqa: E402


# === 终端 ANSI 颜色（对应 DiffView 的 added/removed/normal/meta 着色）===
_ANSI = {
    "added": "\033[32m",    # 绿——新增行
    "removed": "\033[31m",  # 红——删除行
    "normal": "\033[37m",   # 浅灰——未变行
    "meta": "\033[90m",     # 暗灰——元信息
}
_RESET = "\033[0m"
_BOLD = "\033[1m"

# diff 行类型 → 行首前缀符号（即使没有颜色也能区分）
_PREFIX = {
    "added": "+ ",
    "removed": "- ",
    "normal": "  ",
    "meta": "~ ",
}


def _color(tag: str, text: str) -> str:
    """按 tag 给文本套 ANSI 颜色；未知 tag 退回 normal。"""
    code = _ANSI.get(tag, _ANSI["normal"])
    return f"{code}{text}{_RESET}"


# ============================================================
# v1 / v2 文件生成器——每个返回写入磁盘的相对文件名
# ============================================================

def _make_docx(path: Path, *, v2: bool) -> None:
    """报告.docx：v1 两段；v2 改第二段文字 + 新增第三段。"""
    from docx import Document

    doc = Document()
    doc.add_paragraph("第一段：项目背景介绍，本季度目标保持稳定。")
    if not v2:
        doc.add_paragraph("第二段：当前进度 50%，预计下月完成。")
    else:
        # 改一段文字
        doc.add_paragraph("第二段：当前进度 80%，预计本周完成。")
        # 新增一段
        doc.add_paragraph("第三段：新增风险提示，需要额外人力支持。")
    doc.save(str(path))


def _make_pptx(path: Path, *, v2: bool) -> None:
    """演示.pptx：v1 两张；v2 改第一张标题 + 新增第三张。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]  # 纯空白布局，自己加文本框最可控

    def _add_slide(title: str, body: str) -> None:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        tb.text_frame.text = title
        bb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
        bb.text_frame.text = body

    if not v2:
        _add_slide("封面：2026 第一季度汇报", "汇报人：张三")
        _add_slide("第二页：销售概况", "总额 100 万元")
    else:
        # 改第一张标题
        _add_slide("封面：2026 第二季度汇报", "汇报人：张三")
        _add_slide("第二页：销售概况", "总额 100 万元")
        # 新增第三张
        _add_slide("第三页：下季度计划", "拓展华南市场")
    prs.save(str(path))


def _make_xlsx(path: Path, *, v2: bool) -> None:
    """成绩.xlsx：v1 表头 + 两行；v2 改一个单元格值 + 新增一行。"""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "成绩单"
    ws["A1"] = "姓名"
    ws["B1"] = "分数"
    ws["A2"] = "张三"
    ws["A3"] = "李四"
    if not v2:
        ws["B2"] = 85
        ws["B3"] = 90
    else:
        ws["B2"] = 95  # 改一个单元格值（85 → 95）
        ws["B3"] = 90
        # 新增一行
        ws["A4"] = "王五"
        ws["B4"] = 88
    wb.save(str(path))


def _make_pdf(path: Path, *, v2: bool) -> None:
    """手册.pdf：两页；v2 改第一页文字、第二页保持不变。"""
    import fitz  # PyMuPDF

    doc = fitz.open()  # 新建空文档
    page1 = doc.new_page()
    if not v2:
        page1.insert_text((72, 72), "Page 1: Installation guide version 1.0")
    else:
        # 改第一页文字（版本号 1.0 → 2.0）
        page1.insert_text((72, 72), "Page 1: Installation guide version 2.0")
    page2 = doc.new_page()
    # 第二页两个版本完全一致
    page2.insert_text((72, 72), "Page 2: Frequently asked questions")
    doc.save(str(path))
    doc.close()


def _make_png(path: Path, *, v2: bool) -> None:
    """图标.png：v1 纯红 100x100；v2 改尺寸 200x150 + 换成蓝色。"""
    from PIL import Image

    if not v2:
        img = Image.new("RGB", (100, 100), (220, 30, 30))  # 红
    else:
        img = Image.new("RGB", (200, 150), (30, 60, 220))  # 蓝 + 改尺寸
    img.save(str(path), format="PNG")


# 五种格式：(显示名, 文件名, 生成函数)
_SPECS = [
    ("docx", "报告.docx", _make_docx),
    ("pptx", "演示.pptx", _make_pptx),
    ("xlsx", "成绩.xlsx", _make_xlsx),
    ("pdf", "手册.pdf", _make_pdf),
    ("png", "图标.png", _make_png),
]


def _write_version(workspace: Path, *, v2: bool) -> None:
    """把当前版本（v1 或 v2）的 5 个文件全部写到 workspace。"""
    for _label, fname, maker in _SPECS:
        maker(workspace / fname, v2=v2)


def _commit_all(repo: Repository, workspace: Path) -> int:
    """把 workspace 里 5 个文件全部以 'human' 身份 upsert 提交，返回 commit id。"""
    attr = AgentAttribution(agent="human", confidence=1.0)
    changes: list[Change] = []
    for _label, fname, _maker in _SPECS:
        changes.append(
            Change(
                file_path=workspace / fname,
                event_time=0.0,
                attribution=attr,
                kind="upsert",
            )
        )
    # 注意：commit 的签名是 commit(agent, changes, attribution=None)
    return repo.commit(agent="human", changes=changes, attribution=attr)


def run_selftest() -> bool:
    """执行整套自测，返回 True/False 表示是否五种格式都显示出了文字改动。"""
    workspace = Path(tempfile.mkdtemp(prefix="trace_selftest_"))
    print(f"{_BOLD}Trace · 五格式 diff 自测{_RESET}")
    print(f"临时 workspace: {workspace}\n")

    # 记录每种格式：摘要行 + 是否真有 added/removed（即"具体文字改动"）
    summaries: dict[str, str] = {}
    has_change: dict[str, bool] = {}

    try:
        # --- 1. 初始化仓库 + 写 v1 + commit #1 ---
        repo = Repository(workspace)
        repo.init_if_needed()

        _write_version(workspace, v2=False)
        cid1 = _commit_all(repo, workspace)
        print(f"已提交 commit #{cid1}（v1：5 个文件初始版本）")

        # --- 2. 覆盖写 v2 + commit #2 ---
        _write_version(workspace, v2=True)
        cid2 = _commit_all(repo, workspace)
        print(f"已提交 commit #{cid2}（v2：制造文字改动）\n")

        commits = repo.list_commits()
        print(f"仓库现有 commit：{[c['id'] for c in commits]}")
        print("=" * 64)

        # --- 3. 走真实 DiffView 管线：commit#1 → #2 ---
        # 完全复刻 views/diff_view.py：_load_manifest → 比 hash →
        # storage.get → HandlerRegistry.for_path → render_diff
        prev_manifest = repo._load_manifest(cid1)
        cur_manifest = repo._load_manifest(cid2)
        all_paths = sorted(set(prev_manifest.keys()) | set(cur_manifest.keys()))

        # path → 格式 label（按后缀）
        ext_label = {Path(f).suffix.lower(): lbl for lbl, f, _ in _SPECS}

        for path_str in all_paths:
            prev_hash = prev_manifest.get(path_str)
            cur_hash = cur_manifest.get(path_str)
            if prev_hash == cur_hash:
                continue  # 未变化文件跳过（与 DiffView 一致）

            label = ext_label.get(Path(path_str).suffix.lower(), "?")

            # 状态标签（与 DiffView 的 [新增]/[删除]/[修改] 一致）
            if prev_hash is None:
                status = "[新增]"
            elif cur_hash is None:
                status = "[删除]"
            else:
                status = "[修改]"

            print()
            print(f"{_BOLD}{status} {path_str}{_RESET}  （格式：{label}）")

            # 取字节：缺失给空字节，交给 handler 自行处理（同 DiffView）
            old_blob = repo.storage.get(prev_hash) if prev_hash else b""
            new_blob = repo.storage.get(cur_hash) if cur_hash else b""

            # 按文件类型分发到对应 handler
            handler = HandlerRegistry.for_path(Path(path_str))

            # 摘要行（describe_change）
            summary = handler.describe_change(old_blob, new_blob)
            summaries[label] = summary
            print(_color("meta", f"  摘要：{summary}"))

            # 逐行 diff（render_diff）
            diff_lines = handler.render_diff(old_blob, new_blob)
            saw_added_removed = False
            for tag, line in diff_lines:
                prefix = _PREFIX.get(tag, "  ")
                print(_color(tag, f"  {prefix}{line}"))
                if tag in ("added", "removed"):
                    saw_added_removed = True

            # 判定"是否看到具体改动"：
            #   - 文本类（docx/pptx/xlsx/pdf）要求出现 added/removed 行
            #   - 图像类（png）没有逐行文字 diff，以 describe_change 报出
            #     尺寸/大小变化即视为"看到了具体改动"
            if label == "png":
                has_change[label] = "→" in summary and "失败" not in summary
            else:
                has_change[label] = saw_added_removed

        print()
        print("=" * 64)

        # --- 4. 汇总判定 ---
        all_labels = [lbl for lbl, _f, _m in _SPECS]
        ok = all(has_change.get(lbl) for lbl in all_labels)

        print(f"{_BOLD}各格式 handler 摘要：{_RESET}")
        for lbl in all_labels:
            mark = "✅" if has_change.get(lbl) else "❌"
            print(f"  {mark} {lbl:5s} → {summaries.get(lbl, '（无摘要）')}")

        print()
        if ok:
            print(f"{_ANSI['added']}{_BOLD}"
                  f"✅ 自测通过：5 种格式均显示出具体文字改动"
                  f"{_RESET}")
        else:
            missing = [lbl for lbl in all_labels if not has_change.get(lbl)]
            print(f"{_ANSI['removed']}{_BOLD}"
                  f"❌ 自测失败：以下格式未显示出具体改动 → {missing}"
                  f"{_RESET}")
        return ok

    finally:
        # 清理临时目录（无论成功失败都删）
        shutil.rmtree(workspace, ignore_errors=True)


if __name__ == "__main__":
    success = run_selftest()
    sys.exit(0 if success else 1)
